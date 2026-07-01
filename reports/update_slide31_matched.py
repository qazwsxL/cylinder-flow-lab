"""
update_slide31_matched.py
=========================
Once the GPU-node rerender has produced matched-colorbar PNGs in
  reports/figs_v2/pinn_vorticity_matched/
this script:
  1) builds a single 2x2 figure (CFD truth + 3 PINN runs) on a SHARED
     vmin/vmax = ±8 colour scale → reports/figs_v2/fig03_matched_panel.png
  2) replaces the old "auto-colorbar" picture on slide 31 of the v2
     progress deck with this new figure.

Run locally (Mac is fine — no torch needed):

    cd ~/Desktop/cylinder-flow-lab/reports
    python3 update_slide31_matched.py

A backup of the pptx already lives at
  reports/Reconstruction of Flow Fields using PINNs - v2 progress.backup.pptx
(created by add_v2_diagnostic_slides.py).
"""

import os
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.image as mpimg
import pyvista as pv
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
FIGS = os.path.join(ROOT, "reports", "figs_v2")
MATCH = os.path.join(FIGS, "pinn_vorticity_matched")
PPTX  = os.path.join(ROOT, "reports",
                     "Reconstruction of Flow Fields using PINNs - v2 progress.pptx")

# Verify inputs
NEEDED = ["P1_baseline.png", "P2_strict_consist.png", "P2_allCFD_prio2.png"]
for f in NEEDED:
    p = os.path.join(MATCH, f)
    assert os.path.exists(p), f"missing: {p} — did you rsync from Oscar?"
print("[ok] all three matched-colorbar PNGs present")


# ---------------------------------------------------------------------
# 1. CFD truth vorticity, sampled on the same regular grid (-8..12)x(-8..8)
# ---------------------------------------------------------------------
def cfd_vorticity_field():
    m = pv.read(os.path.join(ROOT, "Re40.vtk")).cell_data_to_point_data()
    nx, ny = 400, 200
    xs = np.linspace(-8.0, 12.0, nx)
    ys = np.linspace(-8.0,  8.0, ny)
    grid = pv.ImageData(
        dimensions=(nx, ny, 1),
        spacing=(xs[1] - xs[0], ys[1] - ys[0], 1.0),
        origin=(xs[0], ys[0], 0.0),
    )
    probed = grid.sample(m)
    U = np.asarray(probed.point_data["UMean"]).reshape(ny, nx, 3)
    dx = xs[1] - xs[0]; dy = ys[1] - ys[0]
    W = np.gradient(U[..., 1], dx, axis=1) - np.gradient(U[..., 0], dy, axis=0)
    X, Y = np.meshgrid(xs, ys)
    W[(X**2 + Y**2) < 0.5**2] = np.nan
    return xs, ys, W


xs, ys, W_cfd = cfd_vorticity_field()
print(f"[cfd] vorticity range [{np.nanmin(W_cfd):+.2f}, {np.nanmax(W_cfd):+.2f}]")


# ---------------------------------------------------------------------
# 2. Build the 2x2 figure
# ---------------------------------------------------------------------
fig, axs = plt.subplots(2, 2, figsize=(13.5, 7.4),
                        gridspec_kw=dict(wspace=0.18, hspace=0.30))
axs = axs.flatten()
VLIM = 8.0

# (0) CFD truth — render the field directly so the colorbar matches exactly
im_cfd = axs[0].pcolormesh(xs, ys, W_cfd, shading="auto", cmap="coolwarm",
                           vmin=-VLIM, vmax=VLIM)
theta = np.linspace(0, 2 * np.pi, 80)
axs[0].fill(0.5 * np.cos(theta), 0.5 * np.sin(theta),
            color="0.15", edgecolor="black")
axs[0].set_xlim(-8, 12); axs[0].set_ylim(-8, 8); axs[0].set_aspect("equal")
axs[0].set_title("CFD truth   (|peak ω| ≈ 7.5)", fontsize=11)

# (1..3) matched-colorbar PINN renders — they ALREADY contain their own
# colorbar/labels, so we drop them in as raw images.
LABELS = [
    ("P1 baseline (data-only)",        "P1_baseline.png"),
    ("P2 strict consistency (no data)", "P2_strict_consist.png"),
    ("P2 all-CFD anchor, prio = 2",    "P2_allCFD_prio2.png"),
]
for ax, (title, fname) in zip(axs[1:], LABELS):
    img = mpimg.imread(os.path.join(MATCH, fname))
    ax.imshow(img)
    ax.set_title(title, fontsize=11)
    ax.set_xticks([]); ax.set_yticks([])

# Single colorbar on the right of the CFD panel
plt.colorbar(im_cfd, ax=axs[0], shrink=0.85, label="ω")

fig.suptitle("Vorticity — CFD truth vs three small-net (32/3) PINN runs   "
             "[same domain, same colorbar  vmin/vmax = ±8]",
             y=1.02, fontsize=12)

out_path = os.path.join(FIGS, "fig03_matched_panel.png")
fig.tight_layout()
fig.savefig(out_path, dpi=140, bbox_inches="tight")
plt.close(fig)
print(f"[fig] wrote {out_path}")


# ---------------------------------------------------------------------
# 3. Replace the big picture on slide 31 (0-indexed 30) of the pptx
# ---------------------------------------------------------------------
prs = Presentation(PPTX)
sl = prs.slides[30]   # slide 31

removed = 0
for sh in list(sl.shapes):
    # The old vorticity figure is the only LEFT-half picture (left < 5in,
    # width > 7in).  This avoids killing the small text-box icons on the
    # right side cards.
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE \
       and sh.left < Inches(5) and sh.width > Inches(7):
        sp = sh._element; sp.getparent().remove(sp)
        removed += 1
print(f"[pptx] removed {removed} old picture(s) from slide 31")

# Drop in the new figure at the same anchor the old picture used:
# left=0.45, top=1.30, max_w=8.50, max_h=5.50  (see add_v2_diagnostic_slides.py)
def add_picture_fit(slide, image_path, left, top, max_w, max_h):
    from PIL import Image
    im = Image.open(image_path)
    iw, ih = im.size
    a = iw / ih; ta = max_w / max_h
    if a > ta:
        w = max_w; h = max_w / a
        x = left;          y = top + (max_h - h) / 2
    else:
        h = max_h; w = max_h * a
        y = top;           x = left + (max_w - w) / 2
    slide.shapes.add_picture(image_path, Inches(x), Inches(y), Inches(w), Inches(h))

add_picture_fit(sl, out_path, left=0.45, top=1.30, max_w=8.50, max_h=5.50)
print("[pptx] inserted fig03_matched_panel.png at (0.45, 1.30)  max 8.50 x 5.50 in")

# Optionally update the subtitle line on slide 31 — the old one said the
# colorbars were auto-scaled.  Now they match, so let's rewrite that string.
for sh in sl.shapes:
    if not sh.has_text_frame:
        continue
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            if "auto-scaled" in r.text or "rerender script provided separately" in r.text:
                r.text = ("Same small net 32/3.  All four panels now share the "
                          "vmin/vmax = ±8 colorbar (matched rerender of the PINN "
                          "checkpoints on Oscar).")
                print("[pptx] subtitle updated")
                break

prs.save(PPTX)
print(f"[pptx] saved -> {PPTX}")
print("\nDone.  Re-render the slide and visually QA if you want to double-check.")
