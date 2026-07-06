"""
add_v2_diagnostic_slides.py
===========================
Appends four new slides to the v2 progress deck covering this iteration's
diagnostics (per mentor's feedback):

  Slide 28 — Phase-1 baseline check (P1 mae_u not actually "small" on 32/3)
  Slide 29 — Loss-vs-iteration curves for the three runs
  Slide 30 — Velocity-error comparison, using √(u²+v²) normalization
             instead of the misleading rel_l2_v
  Slide 31 — Vorticity: CFD truth vs three PINN runs

Style mirrors slide 23/27 (BLANK layout, navy/red palette, Calibri/Consolas).
"""

import os, shutil
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT  = "/sessions/dazzling-great-fermi/mnt/cylinder-flow-lab"
SRC   = os.path.join(ROOT, "reports",
                     "Reconstruction of Flow Fields using PINNs - v2 progress.pptx")
DST   = os.path.join(ROOT, "reports",
                     "Reconstruction of Flow Fields using PINNs - v2 progress.pptx")
FIGS  = os.path.join(ROOT, "reports", "figs_v2")

# Restore from backup so this script is idempotent.
BACKUP = os.path.join(ROOT, "reports",
                      "Reconstruction of Flow Fields using PINNs - v2 progress.backup.pptx")
if os.path.exists(BACKUP):
    shutil.copy(BACKUP, SRC)
    print(f"restored {SRC} from backup")
else:
    shutil.copy(SRC, BACKUP)
    print(f"backup -> {BACKUP}")

# Palette taken from slide 23 / 27 inspection
DARK   = RGBColor(0x21, 0x21, 0x21)
GRAY   = RGBColor(0x4A, 0x4A, 0x4A)
NAVY   = RGBColor(0x1E, 0x27, 0x61)
RED    = RGBColor(0x7A, 0x1B, 0x1B)
GREEN  = RGBColor(0x2C, 0x5F, 0x2D)
LIGHT  = RGBColor(0xEC, 0xEF, 0xF4)


def add_text(slide, left, top, width, height, text_blocks, *, anchor=MSO_ANCHOR.TOP):
    """text_blocks = list of (text, dict_of_run_props).  Each block becomes a paragraph."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    for i, (text, props) in enumerate(text_blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        f = run.font
        f.size = Pt(props.get("size", 11))
        f.bold = props.get("bold", False)
        f.italic = props.get("italic", False)
        f.name = props.get("name", "Calibri")
        col = props.get("color", DARK)
        f.color.rgb = col
        if "align" in props:
            p.alignment = props["align"]
    return box


def add_card(slide, left, top, width, height, fill=LIGHT, outline=GRAY):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    rect.fill.solid(); rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = outline
    rect.line.width = Pt(0.75)
    # Make sure card has no text overlay
    rect.text_frame.clear()
    return rect


def add_title(slide, title, subtitle, slide_no):
    # top accent strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   0, 0, Emu(12192000), Inches(0.18))
    strip.fill.solid(); strip.fill.fore_color.rgb = NAVY
    strip.line.fill.background()

    add_text(slide, 0.55, 0.28, 12.5, 0.52,
             [(title, {"size": 26, "bold": True, "color": DARK})])
    add_text(slide, 0.57, 0.85, 12.5, 0.30,
             [(subtitle, {"size": 11, "color": GRAY})])

    line = slide.shapes.add_connector(1, Inches(0.55), Inches(1.18),
                                       Inches(12.78), Inches(1.18))
    line.line.color.rgb = GRAY
    line.line.width = Pt(0.5)

    # footer
    add_text(slide, 0.55, 7.05, 8.0, 0.30,
             [("Junhe Chen  |  PINN Cylinder Flow Reconstruction  |  Brown University",
               {"size": 9, "color": GRAY})])
    add_text(slide, 12.5, 7.05, 0.7, 0.30,
             [(str(slide_no), {"size": 9, "color": GRAY})])


def add_picture_fit(slide, image_path, left, top, max_w, max_h):
    """Add an image, scaled to fit within max_w × max_h while preserving aspect."""
    from PIL import Image
    im = Image.open(image_path)
    iw, ih = im.size
    aspect = iw / ih
    target_aspect = max_w / max_h
    if aspect > target_aspect:
        w = max_w
        h = max_w / aspect
        x = left
        y = top + (max_h - h) / 2
    else:
        h = max_h
        w = max_h * aspect
        y = top
        x = left + (max_w - w) / 2
    slide.shapes.add_picture(image_path, Inches(x), Inches(y),
                             Inches(w), Inches(h))


# =====================================================================
# Open and append
# =====================================================================
prs = Presentation(SRC)
BLANK = prs.slide_layouts[0]
N_existing = len(prs.slides)
print(f"existing slides: {N_existing}")


# ---------------------------------------------------------------------
# SLIDE 28 — Phase-1 baseline isn't small enough on small net
# ---------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK)
add_title(sl,
          "Phase-1 Baseline — How Small Is Small?",
          "Mentor's check: P1 training loss should be very small.  "
          "On the 32/3 net it isn't — mae_u plateaus at 0.106.",
          N_existing + 1)

# Left card — 32/3 numbers
add_card(sl, 0.55, 1.35, 6.10, 2.55, fill=RGBColor(0xF8, 0xEE, 0xEC), outline=RED)
add_text(sl, 0.72, 1.45, 5.85, 0.30,
         [("Small net  (width = 32, depth = 3, 2.8 k params)",
           {"size": 13, "bold": True, "color": RED})])
add_text(sl, 0.72, 1.80, 5.85, 2.10, [
    ("training-time DATA MSE (BFGS plateau)",  {"size": 10.5, "bold": True, "color": DARK}),
    ("    0.036    (~ rmse 0.19)",             {"size": 10.5, "color": NAVY, "name": "Consolas"}),
    ("full-mesh mae_u (7456 cells)",            {"size": 10.5, "bold": True, "color": DARK}),
    ("    0.106",                               {"size": 10.5, "color": NAVY, "name": "Consolas"}),
    ("full-mesh mae_v",                         {"size": 10.5, "bold": True, "color": DARK}),
    ("    0.068",                               {"size": 10.5, "color": NAVY, "name": "Consolas"}),
    ("normalized speed-vector error",           {"size": 10.5, "bold": True, "color": DARK}),
    ("    mean|du,dv| / mean_speed  =  17.0 %", {"size": 10.5, "color": NAVY, "name": "Consolas"}),
])

# Right card — 64/4 numbers
add_card(sl, 6.85, 1.35, 6.10, 2.55, fill=RGBColor(0xEC, 0xF4, 0xEE), outline=GREEN)
add_text(sl, 7.02, 1.45, 5.85, 0.30,
         [("Scale-up net  (width = 64, depth = 4, 14.9 k params)",
           {"size": 13, "bold": True, "color": GREEN})])
add_text(sl, 7.02, 1.80, 5.85, 2.10, [
    ("training-time DATA MSE (BFGS plateau)",  {"size": 10.5, "bold": True, "color": DARK}),
    ("    ~ 0.0008    (~ rmse 0.028)",          {"size": 10.5, "color": NAVY, "name": "Consolas"}),
    ("full-mesh mae_u",                         {"size": 10.5, "bold": True, "color": DARK}),
    ("    0.0186          (5.7× better)",       {"size": 10.5, "color": NAVY, "name": "Consolas"}),
    ("full-mesh mae_v",                         {"size": 10.5, "bold": True, "color": DARK}),
    ("    0.0108",                              {"size": 10.5, "color": NAVY, "name": "Consolas"}),
    ("normalized speed-vector error",           {"size": 10.5, "bold": True, "color": DARK}),
    ("    mean|du,dv| / mean_speed  =  2.9 %",  {"size": 10.5, "color": NAVY, "name": "Consolas"}),
])

# Bottom finding card
add_card(sl, 0.55, 4.10, 12.40, 2.45,
         fill=RGBColor(0xFB, 0xF8, 0xE9), outline=RGBColor(0xC0, 0xA8, 0x44))
add_text(sl, 0.75, 4.20, 12.10, 2.25, [
    ("Finding",
        {"size": 13, "bold": True, "color": RED}),
    ("•  P1 \"should be small\" only holds with sufficient capacity.  "
     "On 32/3 the BFGS plateau (0.036) and full-mesh mae_u (0.106) reveal a "
     "representation ceiling, not an optimization failure.",
        {"size": 11.5, "color": DARK}),
    ("•  This is exactly the bottleneck that slide 27's scale-up to 64/4 "
     "breaks: same protocol, same data_priority, mae_u drops by 5.7× and "
     "normalized speed error from 17.0% → 2.9%.",
        {"size": 11.5, "color": DARK}),
    ("•  All subsequent diagnostic comparisons (loss curves, vorticity) "
     "use the 32/3 baseline because the failure modes are clearer there.",
        {"size": 11.5, "color": DARK}),
])


# ---------------------------------------------------------------------
# SLIDE 29 — Loss curves for the three runs
# ---------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK)
add_title(sl,
          "Loss Curves — Three Runs Side by Side",
          "P1 baseline (data-only)   /   P2 strict consistency (no data)   /   "
          "P2 all-CFD anchor, prio = 2.0     [small net 32/3, Adam ⟶ BFGS on a single x-axis]",
          N_existing + 2)

add_picture_fit(sl, os.path.join(FIGS, "fig01_loss_curves.png"),
                left=0.45, top=1.25, max_w=12.45, max_h=2.75)
add_picture_fit(sl, os.path.join(FIGS, "fig01b_mae_curves.png"),
                left=0.45, top=4.05, max_w=12.45, max_h=2.30)

add_text(sl, 0.55, 6.40, 12.40, 0.55, [
    ("Top: training-loss components that are actually in each run's objective  "
     "(DATA off for strict-consistency, PDE off for baseline).",
     {"size": 9.5, "color": GRAY, "italic": True}),
    ("Bottom: full-mesh CFD-monitor MAE — note strict-consistency drifts UP from 0.10 to 0.34 once data anchor is removed.",
     {"size": 9.5, "color": GRAY, "italic": True}),
])


# ---------------------------------------------------------------------
# SLIDE 30 — Velocity-error comparison with √(u²+v²) normalization
# ---------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK)
add_title(sl,
          "Velocity Error — √(u²+v²) Normalization, not rel_l2_v",
          "Mentor's correction: ‖v_true‖₂ is only 11.5 % of ‖u_true‖₂, "
          "so rel_l2_v inflates the v error by ≈ 8.7×.  Use speed-magnitude normalization instead.",
          N_existing + 3)

# Bar chart on the left
add_picture_fit(sl, os.path.join(FIGS, "fig02_metric_bars.png"),
                left=0.40, top=1.30, max_w=8.45, max_h=4.50)

# Table on the right
add_card(sl, 9.05, 1.35, 3.95, 4.45,
         fill=RGBColor(0xF6, 0xF7, 0xFB), outline=NAVY)
add_text(sl, 9.15, 1.45, 3.85, 0.30,
         [("Final-iteration metrics", {"size": 13, "bold": True, "color": NAVY})])
add_text(sl, 9.15, 1.75, 3.85, 4.00, [
    ("CFD truth (7456 cells)",                     {"size": 10, "bold": True, "color": DARK}),
    ("    ‖u‖₂ = 77.17                            ", {"size": 9.5,  "color": NAVY, "name": "Consolas"}),
    ("    ‖v‖₂ =  8.91     (only 11.5% of ‖u‖₂)",  {"size": 9.5,  "color": NAVY, "name": "Consolas"}),
    ("    mean_speed = 0.816",                      {"size": 9.5,  "color": NAVY, "name": "Consolas"}),
    ("P1 baseline (data-only)",                    {"size": 10, "bold": True, "color": DARK}),
    ("    mae_u = 0.106   mae_v = 0.068",          {"size": 9.5,  "color": NAVY, "name": "Consolas"}),
    ("    speed-norm  =  17.0 %    rel_l2 =  0.92", {"size": 9.5,  "color": NAVY, "name": "Consolas"}),
    ("P2 strict consistency (no data)",            {"size": 10, "bold": True, "color": RED}),
    ("    mae_u = 0.339   mae_v = 0.183",          {"size": 9.5,  "color": RED, "name": "Consolas"}),
    ("    speed-norm  =  49.4 %    rel_l2 =  2.25", {"size": 9.5,  "color": RED, "name": "Consolas"}),
    ("P2 all-CFD anchor (prio = 2)",                {"size": 10, "bold": True, "color": GREEN}),
    ("    mae_u = 0.129   mae_v = 0.071",          {"size": 9.5,  "color": NAVY, "name": "Consolas"}),
    ("    speed-norm  =  19.6 %    rel_l2 =  0.97", {"size": 9.5,  "color": NAVY, "name": "Consolas"}),
])

# Reading card at bottom
add_card(sl, 0.55, 5.95, 12.40, 1.00,
         fill=RGBColor(0xFB, 0xF8, 0xE9), outline=RGBColor(0xC0, 0xA8, 0x44))
add_text(sl, 0.75, 6.05, 12.10, 0.85, [
    ("Reading",      {"size": 12, "bold": True, "color": RED}),
    ("Speed-normalized error makes the three runs comparable on a 0–1 scale:  17 % / 49 % / 20 %.  "
     "Strict-consistency's ≈ 50 % drift is the trivial-attractor failure.  "
     "Old rel_l2_v ≈ 0.9 / 2.2 / 1.0 looked alarming but only because ‖v_true‖₂ is tiny — the numbers were not directly comparable across runs.",
        {"size": 10.5, "color": DARK}),
])


# ---------------------------------------------------------------------
# SLIDE 31 — Vorticity panel
# ---------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK)
add_title(sl,
          "Vorticity — CFD truth vs Three PINN Runs",
          "Same small net 32/3.  PINN colorbars are auto-scaled per panel "
          "(matched-colorbar rerender script provided separately).",
          N_existing + 4)

add_picture_fit(sl, os.path.join(FIGS, "fig03_vorticity_panel.png"),
                left=0.45, top=1.30, max_w=8.50, max_h=5.50)

# Right-hand explanation cards
add_card(sl, 9.10, 1.30, 3.85, 1.55,
         fill=RGBColor(0xF6, 0xF7, 0xFB), outline=NAVY)
add_text(sl, 9.25, 1.40, 3.65, 0.30,
         [("CFD truth", {"size": 12, "bold": True, "color": NAVY})])
add_text(sl, 9.25, 1.70, 3.65, 1.10, [
    ("Two opposite-signed shear-layer bands, |ω|_peak ≈ 7.5; clean recirculation bubble.",
        {"size": 10, "color": DARK}),
])

add_card(sl, 9.10, 2.95, 3.85, 1.40,
         fill=RGBColor(0xF6, 0xF7, 0xFB), outline=NAVY)
add_text(sl, 9.25, 3.05, 3.65, 0.30,
         [("P1 baseline (blue)", {"size": 12, "bold": True, "color": NAVY})])
add_text(sl, 9.25, 3.35, 3.65, 1.00, [
    ("Vorticity is essentially 0 everywhere (auto-range ± 0.06).  "
     "Small net cannot represent the wake's high gradients — velocity looks fine because it is smooth.",
        {"size": 10, "color": DARK}),
])

add_card(sl, 9.10, 4.45, 3.85, 1.40,
         fill=RGBColor(0xF8, 0xEE, 0xEC), outline=RED)
add_text(sl, 9.25, 4.55, 3.65, 0.30,
         [("P2 strict consistency (red)", {"size": 12, "bold": True, "color": RED})])
add_text(sl, 9.25, 4.85, 3.65, 1.00, [
    ("Vorticity ± 7.5 noise spread across the whole domain — no coherent wake.  "
     "PDE residual → 0 is achieved by drifting to a trivial NS solution.",
        {"size": 10, "color": DARK}),
])

add_card(sl, 9.10, 5.95, 3.85, 1.05,
         fill=RGBColor(0xEC, 0xF4, 0xEE), outline=GREEN)
add_text(sl, 9.25, 6.00, 3.65, 0.30,
         [("P2 all-CFD anchor (green)", {"size": 12, "bold": True, "color": GREEN})])
add_text(sl, 9.25, 6.30, 3.65, 0.70, [
    ("Vorticity localised near cylinder, qualitatively closer to CFD; "
     "needs scale-up for derivative-level match.",
        {"size": 10, "color": DARK}),
])


# =====================================================================
# Save
# =====================================================================
prs.save(DST)
print(f"saved -> {DST}")
print(f"new slide count: {len(prs.slides)}")
